import io
import subprocess
import tempfile
import os
from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse

router = APIRouter()


@router.post("/pdf-to-ppt")
async def pdf_to_ppt(
    files: list[UploadFile] = File(...),
    options: str = Form(default="{}"),
):
    """Convert PDF to PPTX (each page becomes a slide with image)."""
    if not files:
        raise HTTPException(status_code=400, detail="請上傳 PDF 檔案")

    try:
        from pdf2image import convert_from_bytes
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise HTTPException(status_code=501, detail="需要安裝 pdf2image、python-pptx 和 poppler")

    data = await files[0].read()
    images = convert_from_bytes(data, dpi=150)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        buf = io.BytesIO()
        img.save(buf, "JPEG", quality=85)
        buf.seek(0)
        slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=converted.pptx"},
    )
